import os
import re
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ──────────────────────────────────────────────────────────────────────────────
# Functions for creating simple table PPT (from create_table.py)
# ──────────────────────────────────────────────────────────────────────────────
def create_ppt_with_table(output_file):
    try:
        # Create a new presentation
        prs = Presentation()

        # Add a blank slide layout
        slide_layout = prs.slide_layouts[5]  # 5 = blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Define table position and size
        rows, cols = 4, 3
        left = Inches(1.0)
        top = Inches(1.5)
        width = Inches(6.0)
        height = Inches(1.5)

        # Add table to slide
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Set column widths (fixed to avoid float issues)
        width_emus = width.emus
        col_width_emus = width_emus // cols
        remainder = width_emus % cols
        for col in range(cols):
            table.columns[col].width = Emu(col_width_emus + (1 if col < remainder else 0))

        # Fill table data
        data = [
            ["Name", "Age", "City"],
            ["Alice", "30", "New York"],
            ["Bob", "25", "Los Angeles"],
            ["Charlie", "35", "Chicago"]
        ]

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = data[r][c]

                # Format text
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Arial"
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                paragraph.alignment = PP_ALIGN.CENTER

        # Save the presentation
        prs.save(output_file)
        print(f"PowerPoint saved as '{output_file}'")

    except Exception as e:
        print(f"Error creating PowerPoint: {e}")

# ──────────────────────────────────────────────────────────────────────────────
# Functions for creating tables from Markdown (from create_tables.py, with fixes)
# ──────────────────────────────────────────────────────────────────────────────
def process_markdown_directory(dir_path):
    """
    Processes all .md files in the given directory, creating a corresponding _tables.pptx file for each.
    """
    for filename in os.listdir(dir_path):
        if filename.endswith('.md'):
            md_path = os.path.join(dir_path, filename)
            ppt_filename = filename.replace('.md', '_tables.pptx')
            ppt_path = os.path.join(dir_path, ppt_filename)
            create_ppt_from_md(md_path, ppt_path)

def create_ppt_from_md(md_path, ppt_path):
    """
    Parses the markdown file to find slide headers and tables, then creates a PowerPoint file with slides for each table.
    """
    try:
        with open(md_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        prs = Presentation()
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if line.startswith('#### Slide '):
                # Extract the title after '#### '
                title = line[5:].strip()

                # Skip lines until the start of the table
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('|'):
                    i += 1

                if i >= len(lines):
                    break

                # Collect table lines
                table_lines = []
                while i < len(lines) and lines[i].strip().startswith('|'):
                    table_lines.append(lines[i].strip())
                    i += 1

                # Parse the table if valid
                if len(table_lines) < 2:
                    continue  # Not a valid table

                # Extract headers (first row)
                headers = [h.strip() for h in table_lines[0].split('|')[1:-1]]

                # Extract data rows (skip separator row)
                data_rows = []
                for row_line in table_lines[2:]:
                    row = [c.strip() for c in row_line.split('|')[1:-1]]
                    if len(row) != len(headers):
                        continue  # Skip uneven rows
                    data_rows.append(row)

                cols = len(headers)
                rows = 1 + len(data_rows)  # Header row + data rows
                if rows < 2 or cols < 1:
                    continue

                # Use blank layout for more control
                slide_layout = prs.slide_layouts[6]  # Blank (fallback to 5 if needed)
                slide = prs.slides.add_slide(slide_layout)

                # Add title as textbox
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                tf = title_box.text_frame
                tf.text = title
                p = tf.paragraphs[0]
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER

                # Define table position and size
                left = Inches(1.0)
                top = Inches(2.0)
                width = Inches(8.0)
                height = Inches(0.6 * rows)

                # Add table to slide
                table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                table = table_shape.table

                # Set column widths equally (fixed to avoid float error)
                width_emus = width.emus
                col_width_emus = width_emus // cols
                remainder = width_emus % cols
                for col in range(cols):
                    table.columns[col].width = Emu(col_width_emus + (1 if col < remainder else 0))

                # Fill header row
                for c in range(cols):
                    cell = table.cell(0, c)
                    cell.text = headers[c]
                    format_cell(cell)

                # Fill data rows
                for r in range(len(data_rows)):
                    for c in range(cols):
                        cell = table.cell(r + 1, c)
                        cell.text = data_rows[r][c]
                        format_cell(cell)

            else:
                i += 1

        # Save the presentation if slides were added
        if len(prs.slides) > 0:
            prs.save(ppt_path)
            print(f"PowerPoint saved as '{ppt_path}'")
        else:
            print(f"No tables found in '{md_path}'. No PPT created.")

    except Exception as e:
        print(f"Error processing '{md_path}': {e}")

def format_cell(cell):
    """
    Formats the text in a table cell.
    """
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.name = "Arial"
    paragraph.font.color.rgb = RGBColor(0, 0, 0)
    paragraph.alignment = PP_ALIGN.CENTER

# ──────────────────────────────────────────────────────────────────────────────
# Functions for creating PPT from text (from extract_pptx.py)
# ──────────────────────────────────────────────────────────────────────────────
def create_slide_with_title_and_bullets(prs, title_text, bullets):
    """
    Add a new slide with title and bulleted content.
    bullets: list of strings like "  - Sub bullet" (indentation determines level)
    """
    # Use Title and Content layout (index 1 in most templates)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_placeholder = slide.shapes.title
    if title_text and title_text != "(No title)":
        title_placeholder.text = title_text
    else:
        title_placeholder.text = "Slide"  # Fallback
    
    # Get the content placeholder (body)
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear any default text (like "Click to add text")
    
    for bullet_line in bullets:
        # Extract indentation and text
        match = re.match(r"^(\s*)- (.*)$", bullet_line)
        if not match:
            continue  # Skip invalid lines
        
        indent_spaces, text = match.groups()
        text = text.strip()
        if not text:
            continue
        
        # Calculate level: 2 spaces = level 1, 4 = level 2, etc.
        level = len(indent_spaces) // 2
        
        p = tf.add_paragraph()
        p.text = text
        p.level = min(level, 8)  # PowerPoint supports up to level 8
        p.font.size = Pt(18)  # Optional: adjust font size

def parse_text_file(txt_path):
    """
    Parse the extracted text file into a list of (title, bullets) per slide.
    Returns: list of tuples (title_str, list_of_bullet_strings)
    """
    if not os.path.exists(txt_path):
        raise FileNotFoundError(f"Text file not found: {txt_path}")
    
    with open(txt_path, "r", encoding="utf-8") as f:
        lines = [line.rstrip() for line in f.readlines()]
    
    slides_data = []
    current_title = None
    current_bullets = []
    in_body_section = False

    for line in lines:
        line = line.rstrip()
        
        if line.startswith("Slide "):
            if current_title is not None:
                slides_data.append((current_title, current_bullets))
            current_title = None
            current_bullets = []
            in_body_section = False
        
        elif line.startswith("Title: "):
            current_title = line[7:].strip()
            if current_title == "(No title)":
                current_title = ""
        
        elif line.strip() == "Body Bullets:":
            in_body_section = True
        
        elif in_body_section and line.strip().startswith("- "):
            # Preserve original indentation (important!)
            current_bullets.append(line)  # Keep as-is, including leading spaces
        
        elif line.strip() == "" or line.startswith("="):
            continue  # Skip empty lines and separators
    
    # Don't forget the last slide
    if current_title is not None:
        slides_data.append((current_title, current_bullets))
    
    return slides_data

def create_pptx_from_text(txt_path, output_pptx=None):
    prs = Presentation()
    
    # Optional: set a nicer default template (minimalist)
    # You can load a custom .potx template here if desired
    
    slides_data = parse_text_file(txt_path)
    
    for title, bullets in slides_data:
        create_slide_with_title_and_bullets(prs, title, bullets)
    
    # Default output filename
    if not output_pptx:
        base = os.path.splitext(os.path.basename(txt_path))[0]
        output_pptx = f"{base}_recreated.pptx"
    
    prs.save(output_pptx)
    print(f"PowerPoint file created: {output_pptx}")
    print(f"   Created {len(slides_data)} slides.")

# ──────────────────────────────────────────────────────────────────────────────
# New function: Extract notes/text from PPTX to structured text file
# ──────────────────────────────────────────────────────────────────────────────
def extract_to_text(pptx_path, output_txt):
    """
    Extracts slide titles and bulleted content from a PPTX file into a structured text file.
    Format similar to example_input.txt.
    """
    try:
        prs = Presentation(pptx_path)
        with open(output_txt, 'w', encoding='utf-8') as f:
            for slide_num, slide in enumerate(prs.slides, start=1):
                f.write(f"Slide {slide_num}\n")
                f.write("====================\n")

                # Extract title
                title = "(No title)"
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
                f.write(f"Title: {title}\n\n")

                # Extract body bullets
                f.write("Body Bullets:\n")
                bullets_found = False
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    tf = shape.text_frame
                    # Skip if this is the title shape
                    if shape == slide.shapes.title:
                        continue
                    for paragraph in tf.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            indent = '  ' * paragraph.level
                            f.write(f"{indent}- {text}\n")
                            bullets_found = True
                if not bullets_found:
                    f.write("- No bullets\n")
                f.write("\n")

        print(f"Extracted text saved to '{output_txt}'")

    except Exception as e:
        print(f"Error extracting from PowerPoint: {e}")

# ──────────────────────────────────────────────────────────────────────────────
# Main CLI entry point
# ──────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Aggregated PowerPoint utility script. Supports multiple modes for creating and extracting PPTX files."
    )
    subparsers = parser.add_subparsers(dest='mode', required=True, help="Mode of operation")

    # Mode: create_simple_table
    p_simple = subparsers.add_parser('create_simple_table', help="Create a simple hardcoded table PPT")
    p_simple.add_argument('--output', required=True, help="Output PPTX file path")

    # Mode: create_tables (from MD directory)
    p_tables = subparsers.add_parser('create_tables', help="Create table PPTs from Markdown files in a directory")
    p_tables.add_argument('--dir', required=True, help="Directory containing .md files")

    # Mode: create_from_text
    p_text = subparsers.add_parser('create_from_text', help="Create PPT from structured text file (titles and bullets)")
    p_text.add_argument('--input', required=True, help="Input text file path")
    p_text.add_argument('--output', help="Optional output PPTX file path (defaults to input_recreated.pptx)")

    # Mode: extract_to_text
    p_extract = subparsers.add_parser('extract_to_text', help="Extract titles and bullets from PPTX to text file")
    p_extract.add_argument('--input', required=True, help="Input PPTX file path")
    p_extract.add_argument('--output', required=True, help="Output text file path")

    args = parser.parse_args()

    if args.mode == 'create_simple_table':
        create_ppt_with_table(args.output)
    elif args.mode == 'create_tables':
        process_markdown_directory(args.dir)
    elif args.mode == 'create_from_text':
        create_pptx_from_text(args.input, args.output)
    elif args.mode == 'extract_to_text':
        extract_to_text(args.input, args.output)
